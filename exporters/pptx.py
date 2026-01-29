"""
FlowCast PowerPoint Exporter
Creates a beautiful presentation with Books & Balances branding.
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import List, Dict, Optional
import matplotlib.pyplot as plt

# Books & Balances Brand Colors
BRAND_PRIMARY = RGBColor(27, 79, 114)       # Deep navy
BRAND_PRIMARY_LIGHT = RGBColor(46, 134, 171)  # Lighter blue
BRAND_SECONDARY = RGBColor(20, 143, 119)    # Teal
BRAND_POSITIVE = RGBColor(39, 174, 96)      # Green
BRAND_NEGATIVE = RGBColor(231, 76, 60)      # Red
BRAND_TEXT = RGBColor(44, 62, 80)           # Dark text
BRAND_LIGHT = RGBColor(127, 140, 141)       # Light text


def fig_to_image_bytes(fig) -> BytesIO:
    """Convert matplotlib figure to image bytes."""
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    return buf


def create_pptx_report(data: dict, metrics: dict, insights: List[dict],
                       charts: dict, forecast_data: dict = None,
                       currency: str = '£') -> bytes:
    """
    Create a beautiful PowerPoint presentation with Books & Balances branding.

    Args:
        data: Parsed financial data
        metrics: Calculated metrics dictionary
        insights: List of insight dictionaries
        charts: Dictionary of matplotlib figures
        forecast_data: Optional forecast data
        currency: Currency symbol

    Returns:
        PPTX file as bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # ==========================================================================
    # Slide 1: Title Slide
    # ==========================================================================
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient effect (using a shape)
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(3))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BRAND_PRIMARY
    bg_shape.line.fill.background()

    # Brand tag
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.4))
    tag_frame = tag_box.text_frame
    tag_para = tag_frame.paragraphs[0]
    tag_para.text = "BOOKS & BALANCES"
    tag_para.font.size = Pt(12)
    tag_para.font.color.rgb = BRAND_SECONDARY
    tag_para.font.bold = True
    tag_para.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "FlowCast"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Company info
    company = data.get('company', 'Company')
    period = data.get('period', '')

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.5))
    info_frame = info_box.text_frame
    info_para = info_frame.paragraphs[0]
    info_para.text = company
    info_para.font.size = Pt(28)
    info_para.font.bold = True
    info_para.font.color.rgb = BRAND_PRIMARY
    info_para.alignment = PP_ALIGN.CENTER

    period_para = info_frame.add_paragraph()
    period_para.text = period
    period_para.font.size = Pt(18)
    period_para.font.color.rgb = BRAND_LIGHT
    period_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tag2_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tag2_frame = tag2_box.text_frame
    tag2_para = tag2_frame.paragraphs[0]
    tag2_para.text = "Financial clarity, delivered beautifully."
    tag2_para.font.size = Pt(14)
    tag2_para.font.italic = True
    tag2_para.font.color.rgb = BRAND_LIGHT
    tag2_para.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # Slide 2: Financial Health Summary
    # ==========================================================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Financial Health Summary")

    # Key metrics
    key_metrics = [
        ('Gross Margin', metrics.get('gross_margin', {})),
        ('Operating Margin', metrics.get('operating_margin', {})),
        ('Revenue Trend', metrics.get('revenue_trend', {})),
        ('Cost Control', metrics.get('cost_control', {})),
    ]

    # Create metric boxes
    x_positions = [Inches(0.75), Inches(3.75), Inches(6.75), Inches(9.75)]

    for i, (label, metric) in enumerate(key_metrics):
        add_metric_box(
            slide,
            x_positions[i],
            Inches(1.8),
            label,
            metric.get('formatted', 'N/A'),
            metric.get('status', 'neutral')
        )

    # Additional figures
    total_rev = metrics.get('total_revenue', {}).get('value', 0)
    total_profit = metrics.get('total_profit', {}).get('value', 0)
    burn_rate = metrics.get('monthly_burn_rate', {}).get('value', 0)
    largest_exp = metrics.get('largest_expense', {}).get('category', 'N/A')

    figures = [
        ('Total Revenue', f'{currency}{total_rev:,.0f}'),
        ('Operating Profit', f'{currency}{total_profit:,.0f}'),
        ('Monthly Burn Rate', f'{currency}{burn_rate:,.0f}'),
        ('Largest Expense', largest_exp),
    ]

    for i, (label, value) in enumerate(figures):
        add_simple_metric(slide, x_positions[i], Inches(4.2), label, str(value))

    # ==========================================================================
    # Slide 3: Key Insights
    # ==========================================================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Key Insights")

    if insights:
        y_pos = Inches(1.6)
        for insight in insights[:5]:
            add_insight_box(slide, Inches(0.75), y_pos, insight['text'], insight['type'])
            y_pos += Inches(1.0)

    # ==========================================================================
    # Slide 4: Revenue & Profit Trend
    # ==========================================================================
    if 'revenue_profit_trend' in charts:
        slide = prs.slides.add_slide(slide_layout)
        add_slide_title(slide, "Revenue & Profit Trends")

        img_stream = fig_to_image_bytes(charts['revenue_profit_trend'])
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), width=Inches(10))

        # Add caption
        caption_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = "Dashed lines show 6-month forecast with confidence bands"
        caption_para.font.size = Pt(12)
        caption_para.font.italic = True
        caption_para.font.color.rgb = BRAND_LIGHT
        caption_para.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # Slide 5: Operating Profit
    # ==========================================================================
    if 'operating_profit' in charts:
        slide = prs.slides.add_slide(slide_layout)
        add_slide_title(slide, "Operating Profit Analysis")

        img_stream = fig_to_image_bytes(charts['operating_profit'])
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), width=Inches(10))

    # ==========================================================================
    # Slide 6: Profit vs Expenses Trend
    # ==========================================================================
    if 'profit_expenses_trend' in charts:
        slide = prs.slides.add_slide(slide_layout)
        add_slide_title(slide, "Cumulative Profit vs Expenses")

        img_stream = fig_to_image_bytes(charts['profit_expenses_trend'])
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), width=Inches(10))

    # ==========================================================================
    # Slide 7: Administrative Costs Breakdown
    # ==========================================================================
    if 'admin_costs_pie' in charts:
        slide = prs.slides.add_slide(slide_layout)
        add_slide_title(slide, "Administrative Costs Breakdown")

        img_stream = fig_to_image_bytes(charts['admin_costs_pie'])
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.3), width=Inches(10))

    # ==========================================================================
    # Slide 8: Forecast Summary
    # ==========================================================================
    if forecast_data:
        slide = prs.slides.add_slide(slide_layout)
        add_slide_title(slide, "6-Month Forecast")

        breakeven = forecast_data.get('breakeven_analysis', {})
        forecast_months = forecast_data.get('months', [])

        y_pos = Inches(1.8)

        # Period
        if forecast_months:
            add_forecast_item(slide, y_pos, "Forecast Period",
                            f"{forecast_months[0]} to {forecast_months[-1]}")
            y_pos += Inches(0.7)

        # Current status
        currently_profitable = breakeven.get('currently_profitable', False)
        status_text = "Profitable" if currently_profitable else "Operating at Loss"
        add_forecast_item(slide, y_pos, "Current Status", status_text)
        y_pos += Inches(0.7)

        # Projected revenue
        if 'revenue' in forecast_data:
            final_rev = forecast_data['revenue']['values'][-1]
            add_forecast_item(slide, y_pos, "Projected Revenue", f"{currency}{final_rev:,.0f}")
            y_pos += Inches(0.7)

        # Projected profit
        if 'operating_profit' in forecast_data:
            final_op = forecast_data['operating_profit']['values'][-1]
            add_forecast_item(slide, y_pos, "Projected Operating Profit", f"{currency}{final_op:,.0f}")
            y_pos += Inches(0.7)

        # Crossover alert
        crossover = breakeven.get('crossover_month')
        crossover_type = breakeven.get('crossover_type')

        if crossover:
            y_pos += Inches(0.3)
            alert_box = slide.shapes.add_textbox(Inches(0.75), y_pos, Inches(11.5), Inches(0.8))
            alert_frame = alert_box.text_frame
            alert_para = alert_frame.paragraphs[0]

            if crossover_type == 'profit_to_loss':
                alert_para.text = f"⚠️ Warning: At current trends, may turn unprofitable by {crossover}"
                alert_para.font.color.rgb = BRAND_NEGATIVE
            else:
                alert_para.text = f"✨ Outlook: At current trends, may turn profitable by {crossover}"
                alert_para.font.color.rgb = BRAND_POSITIVE

            alert_para.font.size = Pt(18)
            alert_para.font.bold = True
            alert_para.alignment = PP_ALIGN.CENTER

    # ==========================================================================
    # Slide 9: Thank You / Contact
    # ==========================================================================
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BRAND_PRIMARY
    bg_shape.line.fill.background()

    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.text = "Thank You"
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    # Website
    web_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.6))
    web_frame = web_box.text_frame
    web_para = web_frame.paragraphs[0]
    web_para.text = "www.booksandbalances.co.uk"
    web_para.font.size = Pt(20)
    web_para.font.color.rgb = BRAND_SECONDARY
    web_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.5))
    tag_frame = tag_box.text_frame
    tag_para = tag_frame.paragraphs[0]
    tag_para.text = "Financial clarity, delivered beautifully."
    tag_para.font.size = Pt(14)
    tag_para.font.italic = True
    tag_para.font.color.rgb = RGBColor(200, 200, 200)
    tag_para.alignment = PP_ALIGN.CENTER

    # Output to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def add_slide_title(slide, title_text: str):
    """Add a styled title to a slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY


def add_metric_box(slide, x, y, label: str, value: str, status: str):
    """Add a metric box with status indicator."""
    width = Inches(2.75)
    height = Inches(1.8)

    # Status colors
    status_colors = {
        'green': (RGBColor(200, 230, 201), BRAND_POSITIVE),
        'yellow': (RGBColor(255, 249, 196), RGBColor(243, 156, 18)),
        'red': (RGBColor(255, 205, 210), BRAND_NEGATIVE),
        'neutral': (RGBColor(238, 238, 238), BRAND_LIGHT)
    }

    bg_color, indicator_color = status_colors.get(status, status_colors['neutral'])

    # Create box shape
    shape = slide.shapes.add_shape(1, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    # Add label
    label_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.4))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.size = Pt(11)
    label_para.font.color.rgb = BRAND_LIGHT
    label_para.font.bold = True
    label_para.alignment = PP_ALIGN.CENTER

    # Add value
    value_box = slide.shapes.add_textbox(x, y + Inches(0.7), width, Inches(0.8))
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = value
    value_para.font.size = Pt(28)
    value_para.font.bold = True
    value_para.font.color.rgb = indicator_color
    value_para.alignment = PP_ALIGN.CENTER


def add_simple_metric(slide, x, y, label: str, value: str):
    """Add a simple metric without status indicator."""
    # Label
    label_box = slide.shapes.add_textbox(x, y, Inches(2.75), Inches(0.3))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.size = Pt(10)
    label_para.font.color.rgb = BRAND_LIGHT
    label_para.font.bold = True
    label_para.alignment = PP_ALIGN.CENTER

    # Value
    value_box = slide.shapes.add_textbox(x, y + Inches(0.35), Inches(2.75), Inches(0.4))
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = value
    value_para.font.size = Pt(16)
    value_para.font.bold = True
    value_para.font.color.rgb = BRAND_PRIMARY
    value_para.alignment = PP_ALIGN.CENTER


def add_insight_box(slide, x, y, text: str, insight_type: str):
    """Add an insight box with colored left border."""
    # Colors
    bg_colors = {
        'positive': RGBColor(232, 248, 245),
        'negative': RGBColor(253, 237, 236),
        'neutral': RGBColor(235, 245, 251)
    }

    border_colors = {
        'positive': BRAND_POSITIVE,
        'negative': BRAND_NEGATIVE,
        'neutral': BRAND_PRIMARY_LIGHT
    }

    width = Inches(11.5)
    height = Inches(0.8)

    # Main box
    shape = slide.shapes.add_shape(1, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_colors.get(insight_type, bg_colors['neutral'])
    shape.line.fill.background()

    # Left border accent
    border = slide.shapes.add_shape(1, x, y, Inches(0.08), height)
    border.fill.solid()
    border.fill.fore_color.rgb = border_colors.get(insight_type, border_colors['neutral'])
    border.line.fill.background()

    # Icon
    icons = {'positive': '✨', 'negative': '⚠️', 'neutral': '💡'}
    icon = icons.get(insight_type, '💡')

    # Text with icon
    text_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_para = text_frame.paragraphs[0]
    text_para.text = f"{icon}  {text}"
    text_para.font.size = Pt(14)
    text_para.font.color.rgb = BRAND_TEXT


def add_forecast_item(slide, y, label: str, value: str):
    """Add a forecast item row."""
    # Label
    label_box = slide.shapes.add_textbox(Inches(0.75), y, Inches(4), Inches(0.5))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label + ":"
    label_para.font.size = Pt(18)
    label_para.font.color.rgb = BRAND_LIGHT

    # Value
    value_box = slide.shapes.add_textbox(Inches(5), y, Inches(6), Inches(0.5))
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = value
    value_para.font.size = Pt(18)
    value_para.font.bold = True
    value_para.font.color.rgb = BRAND_PRIMARY
